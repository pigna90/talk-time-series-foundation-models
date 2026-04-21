# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "pymupdf>=1.27.2.2",
#     "python-pptx>=1.0.2",
# ]
# ///
"""Convert slides.pdf into slides.pptx with one full-bleed image per slide."""

from pathlib import Path

import pymupdf
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).parent
PDF_PATH = HERE / "slides.pdf"
PPTX_PATH = HERE / "slides.pptx"

# 16:9 at 1920x1080 — matches the reveal deck aspect ratio
SLIDE_W = Emu(13_333_333)
SLIDE_H = Emu(7_500_000)
DPI = 200  # rasterization DPI for each PDF page


def main():
    if not PDF_PATH.exists():
        raise SystemExit(f"PDF not found at {PDF_PATH}")

    doc = pymupdf.open(PDF_PATH)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    tmp_dir = HERE / "assets" / "_pdf_pages"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    zoom = DPI / 72
    mat = pymupdf.Matrix(zoom, zoom)

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = tmp_dir / f"page_{i:03d}.png"
        pix.save(img_path)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H,
        )
        print(f"  slide {i + 1}/{len(doc)}")

    prs.save(PPTX_PATH)
    print(f"\nWrote {PPTX_PATH} ({PPTX_PATH.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
